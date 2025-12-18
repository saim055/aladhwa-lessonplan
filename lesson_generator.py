"""
AI-Powered Lesson Plan Generator
Handles content generation, document creation, and packaging
"""

import os
import json
from datetime import datetime
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
import zipfile
import anthropic
from pathlib import Path

class LessonPlanGenerator:
    def __init__(self):
        self.output_folder = 'output'
        self.template_folder = 'documents'
        os.makedirs(self.output_folder, exist_ok=True)
        
        # Initialize AI client (using environment variable)
        self.ai_client = None  # Will be initialized with API key
    
    def generate_complete_package(self, lesson_data):
        """Generate complete lesson plan package"""
        try:
            print("Step 1: Generating AI content...")
            ai_content = self.generate_ai_content(lesson_data)
            
            print("Step 2: Creating lesson plan document...")
            lesson_doc = self.create_lesson_plan_document(lesson_data, ai_content)
            
            print("Step 3: Creating worksheets...")
            worksheets = self.create_worksheets(lesson_data, ai_content)
            
            print("Step 4: Creating rubrics...")
            rubrics = self.create_rubrics(lesson_data, ai_content)
            
            print("Step 5: Creating question bank...")
            question_bank = self.create_question_bank(lesson_data, ai_content)
            
            print("Step 6: Creating PowerPoint...")
            ppt_file = self.create_powerpoint(lesson_data, ai_content)
            
            print("Step 7: Packaging files...")
            zip_file = self.package_files(lesson_data, [
                lesson_doc, worksheets, rubrics, question_bank, ppt_file
            ])
            
            return {
                'status': 'success',
                'files': {
                    'lesson_plan': lesson_doc,
                    'worksheets': worksheets,
                    'rubrics': rubrics,
                    'question_bank': question_bank,
                    'powerpoint': ppt_file,
                    'package': zip_file
                },
                'download_url': f'/api/download/{os.path.basename(zip_file)}'
            }
        
        except Exception as e:
            print(f"Error in generate_complete_package: {str(e)}")
            import traceback
            traceback.print_exc()
            return {
                'status': 'error',
                'message': str(e)
            }
    
    def generate_ai_content(self, lesson_data):
        """Generate comprehensive lesson content using AI"""
        
        # Build comprehensive prompt
        period_descriptions = {
            1: "introductory/foundational level - students are being introduced to the topic for the first time",
            2: "intermediate level - students have foundational knowledge and are deepening understanding",
            3: "advanced/mastery level - students are applying, synthesizing, and evaluating at highest levels"
        }
        
        period_desc = period_descriptions.get(int(lesson_data['period']), period_descriptions[1])
        
        prompt = f"""You are an expert educational content designer for Al Adhwa Private School in the UAE. Generate a comprehensive, pedagogically-sound lesson plan with the following specifications:

**LESSON DETAILS:**
- Grade: {lesson_data['grade']}
- Subject: {lesson_data['subject']}
- Topic: {lesson_data['topic']}
- Period: {lesson_data['period']} ({period_desc})
- Semester: {lesson_data['semester']}
- Standards: {', '.join(lesson_data['standards']) if lesson_data['standards'] else 'General curriculum standards'}
- Digital Platform: {lesson_data['digital_platform'] if lesson_data['digital_platform'] else 'Not specified'}
- Value: {lesson_data['value']}

**REQUIREMENTS:**

1. **LESSON OBJECTIVES (HOT - Higher Order Thinking):**
   - Create 1 overarching objective that requires analysis, evaluation, or creation
   - Must be appropriate for Period {lesson_data['period']} depth

2. **DIFFERENTIATED OUTCOMES (3 levels, DOK-aligned):**
   - Those needing assistance (DOK 1-2): Recall, understand, basic application
   - Average/Middle ability (DOK 2-3): Application, analysis, problem-solving
   - Upper ability (DOK 3-4): Analysis, evaluation, creation, synthesis
   {"- Gifted/Talented (DOK 4): Advanced synthesis, evaluation, creation with real-world complexity" if lesson_data['gifted_talented'] else ""}

3. **KEY VOCABULARY:**
   - List 8-10 essential terms with brief definitions

4. **RESOURCES REQUIRED:**
   - List specific materials, websites, equipment needed
   {f"- Must include integration with {lesson_data['digital_platform']}" if lesson_data['digital_platform'] else ""}

5. **STARTER (5 minutes):**
   - Create an engaging, real-life scenario that grabs student attention
   - Must connect to UAE context where appropriate
   - Include a thought-provoking question

6. **TEACHING COMPONENT (10 minutes maximum):**
   - Clear, structured content delivery
   - Include specific steps/activities
   {f"- Integrate {lesson_data['digital_platform']} effectively" if lesson_data['digital_platform'] else ""}
   - Focus on key concepts for Period {lesson_data['period']}

7. **COOPERATIVE TASKS (15 minutes):**
   For EACH level, provide:
   - Specific group activity with clear instructions
   - 3-5 DOK-appropriate questions for discussion/completion
   - V/A/K indicators (Visual, Auditory, Kinesthetic)
   
   **Those Needing Assistance (DOK 1-2):**
   - Activity: [detailed description]
   - Questions: [list 3-5 questions]
   - V/A/K: [indicate learning styles]
   
   **Average/Middle (DOK 2-3):**
   - Activity: [detailed description]
   - Questions: [list 3-5 questions]
   - V/A/K: [indicate learning styles]
   
   **Upper Ability (DOK 3-4):**
   - Activity: [detailed description]
   - Questions: [list 3-5 questions]
   - V/A/K: [indicate learning styles]
   
   {"**Gifted/Talented (DOK 4):**\\n   - Activity: [detailed description]\\n   - Questions: [list 3-5 questions]\\n   - V/A/K: [indicate learning styles]" if lesson_data['gifted_talented'] else ""}

8. **INDEPENDENT TASKS (15 minutes):**
   Same structure as Cooperative Tasks, but for individual work
   - Must be different from cooperative tasks
   - Allow for independent demonstration of learning
   
9. **PLENARY (5 minutes):**
   - Impressive conclusion with real-world application
   - Reflection questions
   - Forward-looking connection to next lessons

10. **UAE/ADEK INTEGRATION:**
    - **My Identity** (choose relevant elements from: Arabic Language, History, Heritage, Respect, Compassion, Global Understanding, Belonging, Volunteering, Conservation): How does this lesson connect?
    - **Moral Education** (Character and Morality, Individual and Community, Civic Studies, or Cultural Studies): Specific connection with example
    - **STEAM Integration**: Break down Science, Technology, Engineering, Art, Math connections
    - **Links to Other Subjects**: Specific cross-curricular connections
    - **Environment/Sustainability**: How lesson promotes environmental awareness (if applicable)

11. **SKILLS DEVELOPED:**
    - List 3-4 key skills (e.g., Critical Thinking, Collaboration, Digital Literacy, Problem Solving)

Return the content in a structured JSON format with clear sections."""

        # For now, generate structured content (in production, this would call actual AI API)
        ai_content = self._generate_structured_content(lesson_data, period_desc)
        
        return ai_content
    
    def _generate_structured_content(self, lesson_data, period_desc):
        """Generate structured lesson content"""
        
        topic = lesson_data['topic']
        subject = lesson_data['subject']
        grade = lesson_data['grade']
        period = lesson_data['period']
        
        # This is a comprehensive template that would be filled by AI
        # For demonstration, providing structured template
        
        content = {
            'objectives': f"Students will analyze and evaluate {topic} through investigation, experimentation, and application of {subject} principles to real-world scenarios.",
            
            'differentiated_outcomes': {
                'assistance': f"Identify and describe key characteristics of {topic} with support (DOK 1-2)",
                'average': f"Analyze the relationship between variables in {topic} using data and graphs (DOK 2-3)",
                'upper': f"Evaluate experimental results for {topic}, calculate errors, and justify findings by identifying systematic error sources (DOK 3-4)",
                'gifted': f"Design and conduct an original investigation extending {topic} concepts to novel real-world applications with comprehensive analysis (DOK 4)" if lesson_data['gifted_talented'] else None
            },
            
            'vocabulary': self._generate_vocabulary(topic, subject),
            
            'resources': self._generate_resources(lesson_data),
            
            'starter': self._generate_starter(topic, subject, grade),
            
            'teaching_component': self._generate_teaching_component(lesson_data),
            
            'cooperative_tasks': self._generate_differentiated_tasks(lesson_data, 'cooperative'),
            
            'independent_tasks': self._generate_differentiated_tasks(lesson_data, 'independent'),
            
            'plenary': self._generate_plenary(topic, subject),
            
            'adek_integration': self._generate_adek_integration(lesson_data),
            
            'skills': ["Critical Thinking", "Collaboration", "Digital Literacy", "Problem Solving"]
        }
        
        return content
    
    def _generate_vocabulary(self, topic, subject):
        """Generate vocabulary list"""
        # This would be AI-generated in production
        vocab_templates = {
            'Physics': ['Amplitude', 'Period', 'Frequency', 'Wavelength', 'Velocity', 'Acceleration', 'Force', 'Energy'],
            'Chemistry': ['Molecule', 'Atom', 'Reaction', 'Catalyst', 'Solution', 'Compound', 'Element', 'Bond'],
            'Biology': ['Cell', 'Organism', 'Ecosystem', 'Evolution', 'Genetics', 'Metabolism', 'Homeostasis', 'Species'],
            'Math': ['Variable', 'Equation', 'Function', 'Coefficient', 'Constant', 'Expression', 'Solution', 'Graph'],
            'default': ['Concept', 'Process', 'Analysis', 'Evaluation', 'Application', 'Synthesis', 'Investigation', 'Conclusion']
        }
        
        return vocab_templates.get(subject, vocab_templates['default'])
    
    def _generate_resources(self, lesson_data):
        """Generate resource list"""
        resources = [
            "Laptop/Tablet for each student or group",
            "Calculator (scientific)",
            "Whiteboard and markers",
            "Student notebooks"
        ]
        
        if lesson_data['digital_platform']:
            resources.insert(0, f"{lesson_data['digital_platform']} simulation/platform access")
        
        return resources
    
    def _generate_starter(self, topic, subject, grade):
        """Generate engaging starter"""
        return {
            'activity': f"Real-world connection: Show a 30-second video or image demonstrating {topic} in everyday life (e.g., playground swings, car suspension, building design). Students observe and note 3 things they notice.",
            'question': f"How do you think understanding {topic} helps engineers design safer, more efficient systems in our daily lives?",
            'duration': '5 minutes'
        }
    
    def _generate_teaching_component(self, lesson_data):
        """Generate teaching component"""
        platform = lesson_data['digital_platform'] or "demonstrations and examples"
        
        return {
            'duration': '10 minutes maximum',
            'method': f"Interactive demonstration using {platform}",
            'steps': [
                f"1. Introduce key concept of {lesson_data['topic']}",
                "2. Define independent and dependent variables",
                "3. Demonstrate the relationship using platform/simulation",
                "4. Model data collection process",
                "5. Show how to organize data in tables",
                "6. Demonstrate graphing techniques",
                "7. Guide students in identifying patterns"
            ]
        }
    
    def _generate_differentiated_tasks(self, lesson_data, task_type):
        """Generate differentiated tasks for cooperative or independent work"""
        
        topic = lesson_data['topic']
        tasks = {}
        
        # Those Needing Assistance (DOK 1-2)
        tasks['assistance'] = {
            'activity': f"{'In groups,' if task_type == 'cooperative' else 'Individually,'} students will identify and measure basic characteristics of {topic} using guided worksheets with step-by-step instructions and visual aids.",
            'questions': [
                f"1. What are the main parts/components of {topic}? (DOK 1)",
                f"2. How do you measure [key variable] in this setup? (DOK 1)",
                f"3. Record your measurements in the provided table. (DOK 1)",
                f"4. Describe what happens when you change [one variable]. (DOK 2)",
                f"5. Which setup showed the [largest/smallest] value? Why? (DOK 2)"
            ],
            'vak': 'Visual (diagrams, charts), Kinesthetic (hands-on measurement), Auditory (group discussion)' if task_type == 'cooperative' else 'Visual (worksheets), Kinesthetic (measurements), Auditory (self-explanation)'
        }
        
        # Average/Middle (DOK 2-3)
        tasks['average'] = {
            'activity': f"{'Groups' if task_type == 'cooperative' else 'Students'} will conduct systematic investigations of {topic}, collect data across multiple trials, create graphs, and analyze relationships between variables.",
            'questions': [
                f"1. What pattern do you observe in your data for {topic}? (DOK 2)",
                f"2. Create a graph showing the relationship between [variable A] and [variable B]. (DOK 2)",
                f"3. Explain why the relationship follows this pattern. (DOK 3)",
                f"4. What factors might affect the accuracy of your results? (DOK 3)",
                f"5. Predict what would happen if you doubled [one variable]. Test your prediction. (DOK 3)"
            ],
            'vak': 'Visual (graphs, data visualization), Kinesthetic (experimentation), Auditory (explanation and discussion)' if task_type == 'cooperative' else 'Visual (data analysis), Kinesthetic (investigation), Auditory (verbal reasoning)'
        }
        
        # Upper Ability (DOK 3-4)
        tasks['upper'] = {
            'activity': f"{'Collaborative teams' if task_type == 'cooperative' else 'Students independently'} will design comprehensive investigations, conduct error analysis, evaluate theoretical vs. experimental results, and justify discrepancies through systematic error identification.",
            'questions': [
                f"1. Design an investigation to determine the mathematical relationship in {topic}. (DOK 3)",
                f"2. Analyze your graph's slope and calculate its theoretical value. Compare and explain any differences. (DOK 3)",
                f"3. Calculate the percentage error between your results and theoretical predictions. (DOK 3)",
                f"4. Evaluate your experimental method: What systematic errors exist? How do they affect results? (DOK 4)",
                f"5. Propose improvements to reduce errors and justify why your modifications would work. (DOK 4)"
            ],
            'vak': 'Visual (complex graphs, error bars), Kinesthetic (precise measurements), Auditory (justification and evaluation)' if task_type == 'cooperative' else 'Visual (detailed analysis), Kinesthetic (refined experimentation), Auditory (critical reasoning)'
        }
        
        # Gifted/Talented (DOK 4) - if enabled
        if lesson_data['gifted_talented']:
            tasks['gifted'] = {
                'activity': f"{'Advanced research teams' if task_type == 'cooperative' else 'Individual advanced investigation:'} Design and conduct an original research project extending {topic} to real-world applications. Develop novel experimental setups, conduct comprehensive error analysis, and present findings with professional-level documentation.",
                'questions': [
                    f"1. Design an original investigation applying {topic} principles to solve a real-world engineering problem. (DOK 4)",
                    f"2. Synthesize data from multiple trials using statistical analysis (mean, standard deviation, confidence intervals). (DOK 4)",
                    f"3. Evaluate competing theoretical models and determine which best explains your empirical findings. (DOK 4)",
                    f"4. Create a research proposal for extending this investigation, including hypothesis, methodology, and predicted outcomes. (DOK 4)",
                    f"5. Defend your conclusions: How would you respond to a scientist who challenges your error analysis? (DOK 4)"
                ],
                'vak': 'Visual (professional graphs, presentations), Kinesthetic (advanced experimentation), Auditory (research defense)' if task_type == 'cooperative' else 'Visual (research documentation), Kinesthetic (novel experimental design), Auditory (critical evaluation)'
            }
        
        return tasks
    
    def _generate_plenary(self, topic, subject):
        """Generate impressive plenary"""
        return {
            'duration': '5 minutes',
            'activity': f"Class discussion connecting {topic} to real-world applications",
            'real_world_connection': f"Discuss how understanding {topic} is crucial in modern technology, engineering, and daily life. Show examples from UAE infrastructure, renewable energy projects, or cutting-edge technology.",
            'reflection_questions': [
                f"What was the most surprising thing you learned about {topic} today?",
                f"How could you apply this knowledge to solve a problem in your community?",
                "What questions do you still have that you'd like to investigate further?"
            ],
            'forward_connection': f"Next lesson, we'll explore how {topic} connects to [related advanced concept], building on today's foundations."
        }
    
    def _generate_adek_integration(self, lesson_data):
        """Generate ADEK curriculum integration"""
        return {
            'my_identity': "Relate the scientific principles to UAE's innovation and technological advancement, showing how understanding these concepts contributes to national development and global competitiveness.",
            'moral_education': "Discuss the importance of honesty in scientific reporting, integrity in data collection, and perseverance when experiments don't yield expected results.",
            'steam': {
                'science': f"Investigate {lesson_data['topic']} through experimentation and observation",
                'technology': lesson_data['digital_platform'] if lesson_data['digital_platform'] else "Use digital tools for data collection and analysis",
                'engineering': "Apply concepts to solve real-world engineering challenges",
                'art': "Create visual representations (graphs, diagrams) to communicate findings",
                'math': "Use mathematical formulas, graphing, and statistical analysis"
            },
            'links_to_subjects': "Mathematics (graphing, equations), ICT (digital simulations), Engineering (design thinking), Art (visual communication)",
            'environment': "Discuss how scientific understanding helps develop sustainable technologies and protect our environment"
        }
    
    def create_lesson_plan_document(self, lesson_data, ai_content):
        """Create the filled lesson plan Word document"""
        try:
            # Load template
            template_path = os.path.join(self.template_folder, 'lesson_plan_template.docx')
            doc = Document(template_path)
            
            # Fill in the document
            # This is a simplified version - full implementation would parse and fill each table cell
            self._fill_document_fields(doc, lesson_data, ai_content)
            
            # Save
            filename = f"LessonPlan_{lesson_data['subject']}_{lesson_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.docx"
            output_path = os.path.join(self.output_folder, filename)
            doc.save(output_path)
            
            return output_path
        
        except Exception as e:
            print(f"Error creating lesson plan document: {str(e)}")
            # Create a basic document if template loading fails
            return self._create_basic_lesson_plan(lesson_data, ai_content)
    
    def _fill_document_fields(self, doc, lesson_data, ai_content):
        """Fill document fields with generated content"""
        # This would involve complex table navigation and cell filling
        # Simplified version for demonstration
        
        replacements = {
            'Date:': f"Date: {lesson_data['date']}",
            'SEMESTER:': f"SEMESTER: {lesson_data['semester']}",
            'Grade:': f"Grade: {lesson_data['grade']}",
            'Subject:': f"Subject: {lesson_data['subject']}",
            'Topic:': f"Topic: {lesson_data['topic']}",
            'Period:': f"Period: {lesson_data['period']}",
            'Value:': f"Value: {lesson_data['value']}"
        }
        
        for paragraph in doc.paragraphs:
            for key, value in replacements.items():
                if key in paragraph.text:
                    paragraph.text = paragraph.text.replace(key, value)
    
    def _create_basic_lesson_plan(self, lesson_data, ai_content):
        """Create a basic lesson plan document"""
        doc = Document()
        
        # Add header
        header = doc.add_heading('AL ADHWA PRIVATE SCHOOL LESSON PLAN', 0)
        header.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add basic information
        info_table = doc.add_table(rows=2, cols=4)
        info_cells = info_table.rows[0].cells
        info_cells[0].text = f"Date: {lesson_data['date']}"
        info_cells[1].text = f"Semester: {lesson_data['semester']}"
        info_cells[2].text = f"Grade: {lesson_data['grade']}"
        info_cells[3].text = f"Subject: {lesson_data['subject']}"
        
        info_cells = info_table.rows[1].cells
        info_cells[0].text = f"Topic: {lesson_data['topic']}"
        info_cells[1].text = f"Period: {lesson_data['period']}"
        info_cells[2].text = f"Value: {lesson_data['value']}"
        info_cells[3].text = ""
        
        # Add content sections
        doc.add_heading('Lesson Objectives', 1)
        doc.add_paragraph(ai_content['objectives'])
        
        doc.add_heading('Differentiated Outcomes', 1)
        for level, outcome in ai_content['differentiated_outcomes'].items():
            if outcome:
                doc.add_paragraph(f"{level.capitalize()}: {outcome}", style='List Bullet')
        
        doc.add_heading('Key Vocabulary', 1)
        for word in ai_content['vocabulary']:
            doc.add_paragraph(word, style='List Bullet')
        
        # Save
        filename = f"LessonPlan_{lesson_data['subject']}_{lesson_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.docx"
        output_path = os.path.join(self.output_folder, filename)
        doc.save(output_path)
        
        return output_path
    
    def create_worksheets(self, lesson_data, ai_content):
        """Create differentiated worksheets"""
        doc = Document()
        
        doc.add_heading('DIFFERENTIATED WORKSHEETS', 0)
        doc.add_heading(f'Topic: {lesson_data["topic"]}', 1)
        
        # Create worksheet for each DOK level
        dok_levels = ['DOK Level 1-2', 'DOK Level 2-3', 'DOK Level 3-4']
        if lesson_data['gifted_talented']:
            dok_levels.append('DOK Level 4 (Gifted/Talented)')
        
        for level in dok_levels:
            doc.add_page_break()
            doc.add_heading(f'Worksheet: {level}', 1)
            doc.add_paragraph(f"Name: ________________  Date: ________________")
            doc.add_paragraph(f"Grade: {lesson_data['grade']}  Subject: {lesson_data['subject']}")
            doc.add_paragraph()
            
            # Add relevant tasks
            if 'Level 1-2' in level:
                tasks = ai_content['cooperative_tasks']['assistance']
            elif 'Level 2-3' in level:
                tasks = ai_content['cooperative_tasks']['average']
            elif 'Gifted' in level:
                tasks = ai_content['cooperative_tasks'].get('gifted', ai_content['cooperative_tasks']['upper'])
            else:
                tasks = ai_content['cooperative_tasks']['upper']
            
            doc.add_heading('Activity:', 2)
            doc.add_paragraph(tasks['activity'])
            doc.add_paragraph()
            
            doc.add_heading('Questions:', 2)
            for question in tasks['questions']:
                doc.add_paragraph(question, style='List Number')
                doc.add_paragraph()  # Space for answers
                doc.add_paragraph()
        
        # Save
        filename = f"Worksheets_{lesson_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.docx"
        output_path = os.path.join(self.output_folder, filename)
        doc.save(output_path)
        
        return output_path
    
    def create_rubrics(self, lesson_data, ai_content):
        """Create assessment rubrics"""
        doc = Document()
        
        doc.add_heading('ASSESSMENT RUBRICS', 0)
        doc.add_heading(f'Topic: {lesson_data["topic"]}', 1)
        
        # Create rubric table
        rubric_table = doc.add_table(rows=6, cols=5)
        rubric_table.style = 'Table Grid'
        
        # Header row
        header_cells = rubric_table.rows[0].cells
        header_cells[0].text = 'Criteria'
        header_cells[1].text = 'Excellent (4)'
        header_cells[2].text = 'Proficient (3)'
        header_cells[3].text = 'Developing (2)'
        header_cells[4].text = 'Beginning (1)'
        
        # Criteria rows
        criteria = [
            'Understanding of Concepts',
            'Data Collection',
            'Analysis & Interpretation',
            'Communication',
            'Collaboration'
        ]
        
        for i, criterion in enumerate(criteria, 1):
            cells = rubric_table.rows[i].cells
            cells[0].text = criterion
            cells[1].text = 'Demonstrates exceptional understanding and application'
            cells[2].text = 'Demonstrates solid understanding with minor gaps'
            cells[3].text = 'Demonstrates partial understanding with support needed'
            cells[4].text = 'Demonstrates limited understanding, requires significant support'
        
        # Save
        filename = f"Rubrics_{lesson_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.docx"
        output_path = os.path.join(self.output_folder, filename)
        doc.save(output_path)
        
        return output_path
    
    def create_question_bank(self, lesson_data, ai_content):
        """Create question bank organized by DOK levels"""
        doc = Document()
        
        doc.add_heading('QUESTION BANK', 0)
        doc.add_heading(f'Topic: {lesson_data["topic"]}', 1)
        doc.add_heading(f'Grade: {lesson_data["grade"]} | Subject: {lesson_data["subject"]}', 2)
        
        # Organize questions by DOK level
        dok_sections = {
            'DOK Level 1 (Recall & Reproduction)': [],
            'DOK Level 2 (Skills & Concepts)': [],
            'DOK Level 3 (Strategic Thinking)': [],
            'DOK Level 4 (Extended Thinking)': []
        }
        
        # Extract questions from tasks
        all_tasks = {**ai_content['cooperative_tasks'], **ai_content['independent_tasks']}
        
        for level, tasks in all_tasks.items():
            if isinstance(tasks, dict) and 'questions' in tasks:
                for question in tasks['questions']:
                    # Determine DOK level from question
                    if 'identify' in question.lower() or 'define' in question.lower() or 'list' in question.lower():
                        dok_sections['DOK Level 1 (Recall & Reproduction)'].append(question)
                    elif 'analyze' in question.lower() or 'compare' in question.lower() or 'explain' in question.lower():
                        dok_sections['DOK Level 2 (Skills & Concepts)'].append(question)
                    elif 'evaluate' in question.lower() or 'justify' in question.lower() or 'predict' in question.lower():
                        dok_sections['DOK Level 3 (Strategic Thinking)'].append(question)
                    else:
                        dok_sections['DOK Level 4 (Extended Thinking)'].append(question)
        
        # Add questions to document
        for dok_level, questions in dok_sections.items():
            if questions:
                doc.add_page_break()
                doc.add_heading(dok_level, 1)
                for i, question in enumerate(questions, 1):
                    doc.add_paragraph(f"{i}. {question}", style='List Number')
                    doc.add_paragraph()
        
        # Save
        filename = f"QuestionBank_{lesson_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.docx"
        output_path = os.path.join(self.output_folder, filename)
        doc.save(output_path)
        
        return output_path
    
    def create_powerpoint(self, lesson_data, ai_content):
        """Create PowerPoint presentation"""
        prs = Presentation()
        prs.slide_width = PptInches(10)
        prs.slide_height = PptInches(7.5)
        
        # Title Slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = lesson_data['topic']
        subtitle.text = f"{lesson_data['subject']} | Grade {lesson_data['grade']}\nAl Adhwa Private School"
        
        # Teaching model slides based on selection
        ppt_style = lesson_data.get('ppt_style', '7E Model')
        
        if ppt_style == '7E Model':
            self._create_7e_slides(prs, lesson_data, ai_content)
        elif ppt_style == 'I Do, We Do, You Do':
            self._create_gradual_release_slides(prs, lesson_data, ai_content)
        elif ppt_style == '5E Model':
            self._create_5e_slides(prs, lesson_data, ai_content)
        else:
            self._create_traditional_slides(prs, lesson_data, ai_content)
        
        # Save
        filename = f"Presentation_{lesson_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx"
        output_path = os.path.join(self.output_folder, filename)
        prs.save(output_path)
        
        return output_path
    
    def _create_7e_slides(self, prs, lesson_data, ai_content):
        """Create 7E Model presentation slides"""
        stages = [
            ('Elicit', ai_content['starter']),
            ('Engage', 'Hook students with real-world connection'),
            ('Explore', ai_content['teaching_component']),
            ('Explain', 'Present key concepts and vocabulary'),
            ('Elaborate', ai_content['cooperative_tasks']),
            ('Evaluate', 'Assessment and checking understanding'),
            ('Extend', ai_content['plenary'])
        ]
        
        for stage_name, content in stages:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = f"{stage_name}"
            
            body = slide.placeholders[1]
            tf = body.text_frame
            
            if isinstance(content, dict):
                if 'activity' in content:
                    tf.text = content['activity']
                elif 'method' in content:
                    tf.text = content['method']
                else:
                    tf.text = str(content)
            else:
                tf.text = str(content)
    
    def _create_gradual_release_slides(self, prs, lesson_data, ai_content):
        """Create I Do, We Do, You Do presentation slides"""
        # Similar structure for other models
        pass
    
    def _create_5e_slides(self, prs, lesson_data, ai_content):
        """Create 5E Model presentation slides"""
        pass
    
    def _create_traditional_slides(self, prs, lesson_data, ai_content):
        """Create traditional presentation slides"""
        pass
    
    def package_files(self, lesson_data, file_paths):
        """Package all files into a ZIP"""
        zip_filename = f"LessonPlanPackage_{lesson_data['subject']}_{lesson_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
        zip_path = os.path.join(self.output_folder, zip_filename)
        
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for file_path in file_paths:
                if file_path and os.path.exists(file_path):
                    zipf.write(file_path, os.path.basename(file_path))
        
        return zip_path
